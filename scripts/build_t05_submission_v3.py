from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "T05_과제_제출_보고서_v3_2026-08-29.pptx"
SHOT = ROOT / "docs" / "검증스크린샷" / "2026-08-29T04-21-36Z"
if OUT.exists():
    raise SystemExit(f"refusing to overwrite {OUT}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(55, 57, 53)
CARD = RGBColor(22, 25, 21)
CARD2 = RGBColor(44, 48, 40)
INK = RGBColor(238, 238, 231)
MUTED = RGBColor(178, 181, 166)
LIME = RGBColor(198, 255, 34)
RED = RGBColor(255, 105, 92)
AMBER = RGBColor(255, 181, 50)
BLUE = RGBColor(112, 163, 255)


def tx(slide, text, x, y, w, h, size=16, color=INK, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.1)
    frame.margin_right = Inches(0.1)
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return box


def base(title, kicker="T05 · PULSEBOARD"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.58), Inches(0.22), Inches(2.55), Inches(0.34)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = LIME
    tag.line.fill.background()
    tx(slide, kicker, 0.68, 0.27, 2.35, 0.2, 8, CARD, True)
    tx(slide, title, 0.58, 0.7, 11.9, 0.68, 25, INK, True)
    tx(slide, f"{len(prs.slides):02d}", 12.15, 7.06, 0.6, 0.25, 8, MUTED)
    return slide


def card(slide, x, y, w, h, fill=CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(75, 78, 69)
    return shape


def bullets(items):
    return "\n".join("• " + item for item in items)


def cards(title, left_title, left, right_title, right, kicker="DOCUMENT · SUMMARY"):
    slide = base(title, kicker)
    card(slide, 0.65, 1.55, 5.92, 4.95)
    card(slide, 6.77, 1.55, 5.92, 4.95)
    tx(slide, left_title, 0.9, 1.82, 5.3, 0.35, 13, LIME, True)
    tx(slide, bullets(left), 0.9, 2.35, 5.3, 3.75, 17)
    tx(slide, right_title, 7.02, 1.82, 5.3, 0.35, 13, LIME, True)
    tx(slide, bullets(right), 7.02, 2.35, 5.3, 3.75, 17)


def list_slide(title, items, kicker="DOCUMENT · SUMMARY", size=18):
    slide = base(title, kicker)
    card(slide, 0.65, 1.5, 12.04, 5.05)
    tx(slide, bullets(items), 0.95, 1.88, 11.35, 4.25, size)


def place_image(slide, path, x, y, w, h):
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    slide.shapes.add_picture(
        str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph)
    )


def evidence(title, path, caption, kicker="EVIDENCE · SCREENSHOT"):
    slide = base(title, kicker)
    card(slide, 0.58, 1.42, 12.15, 5.45)
    place_image(slide, path, 0.73, 1.58, 11.85, 5.0)
    tx(slide, caption, 0.78, 6.66, 11.7, 0.25, 10, MUTED)


# 01 Cover
s = base("데이터가 늦어도 화면은 사라지지 않는다", "T05 · 과제 제출 보고서 v3")
tx(s, "실시간 정보판 데이터 신선도 배지", 0.72, 1.45, 7.5, 0.55, 26, INK, True)
tx(s, "fresh · stale · unavailable\n마지막 정상값 유지 + 경과 시간 표시", 0.72, 2.25, 7.2, 1.2, 21, MUTED)
card(s, 0.72, 4.0, 8.1, 1.55)
tx(s, "결과물  realtime-info-board-t05.vercel.app\n소스     github.com/stulss/realtime-info-board-t05", 1.02, 4.38, 7.5, 0.8, 15)
card(s, 9.25, 1.48, 3.0, 4.75, CARD2)
tx(s, "10", 9.7, 2.0, 1.0, 0.7, 36, LIME, True)
tx(s, "고정검사", 10.55, 2.2, 1.3, 0.3, 13, MUTED, True)
tx(s, "5개\n원격 상태", 9.7, 3.15, 2.0, 1.0, 22, INK, True)
tx(s, "32장\n기록 보고서", 9.7, 4.65, 2.1, 1.0, 22, INK, True)

# 02 Overview
cards(
    "OVERVIEW · 프로젝트 개요",
    "문제",
    ["장애가 나면 현재 데이터까지 사라질 수 있음", "마지막 정상 수신 시각을 알기 어려움", "색상만으로 상태를 구분하면 접근성 저하"],
    "해결",
    ["fresh / stale / unavailable 3상태", "장애 시 마지막 정상값·차트 유지", "수신 시각·경과 초·텍스트·aria-label"],
    "docs/01_기획.md",
)

# 03 Sources
cards(
    "SOURCE · 과제 진행 중 기록만 사용",
    "기획·요구",
    ["00 과제 요구사항 매핑", "01 기획", "02 인수인계 문서", "T05 고정검사·상한"],
    "검증·결과",
    ["04 작업 기록", "05 배포", "검증안내서", "AI 3줄·최종 보고서"],
    "DOCUMENT · SOURCE SET",
)

# 04 Goal
cards(
    "GOAL · 기능 하나를 정확히 추가",
    "화면 목표",
    ["우측 상단 신선도 배지", "마지막 정상 수신 시각", "매초 증가하는 ageSeconds", "색상 외 텍스트 상태"],
    "장애 목표",
    ["정상 데이터를 실패 응답으로 덮어쓰지 않음", "마지막 값과 차트 유지", "원인별 stale 문구", "정상값 없을 때 unavailable"],
    "docs/01_기획.md §1",
)

# 05 Stack
cards(
    "TECH STACK · 기존 구성을 유지",
    "애플리케이션",
    ["Next.js App Router", "React · TypeScript strict", "Tailwind CSS 4", "TanStack Query 5"],
    "검증·배포",
    ["Vitest · Testing Library", "Playwright", "Vercel", "새 의존성 추가 없음"],
    "docs/01_기획.md §2",
)

# 06 Principles
cards(
    "DESIGN · 학생이 정한 세 가지 원칙",
    "사용자 경험",
    ["에러 뷰와 성공 뷰를 분리하지 않음", "단일 화면에서 배지만 교체", "Toast·모달 대신 정직한 배지와 문구"],
    "수정 경계",
    ["클라이언트 렌더링 단에서 대응", "외부 API·DB 계층 무수정", "기존 실패 처리 함수 재사용"],
    "docs/01_기획.md §3",
)

# 07 Contract
s = base("CONTRACT · 6값 내부 상태를 3값으로 투영", "docs/01_기획.md §4")
for i, (name, rule, color) in enumerate([
    ("fresh", "최신 정상 응답\nstatus = ok", LIME),
    ("stale", "마지막 정상값 존재\n최신 갱신 실패", AMBER),
    ("unavailable", "정상값이 한 번도 없음\nvalue = null", RED),
]):
    x = 0.72 + i * 4.08
    card(s, x, 1.65, 3.78, 4.55)
    tx(s, name, x + 0.25, 2.0, 3.2, 0.4, 17, color, True)
    tx(s, rule, x + 0.25, 2.85, 3.2, 1.15, 19, INK, True)
    tx(s, "ageSeconds는\n클라이언트가 매초 계산", x + 0.25, 5.05, 3.1, 0.75, 12, MUTED)

# 08 Architecture
s = base("ARCHITECTURE · URL에서 배지까지", "docs/01_기획.md §5")
for i, (head, body) in enumerate([
    ("URL", "simulate\n4개 입력"),
    ("QUERY", "정상 fixture\n+ 실패 API"),
    ("FALLBACK", "이전 값 유지\nlastError 기록"),
    ("PROJECT", "dataStatus\nageSeconds"),
    ("VIEW", "기존 화면\n배지만 교체"),
]):
    x = 0.45 + i * 2.54
    card(s, x, 1.75, 2.25, 4.35)
    tx(s, head, x + 0.2, 2.1, 1.85, 0.3, 12, LIME, True)
    tx(s, body, x + 0.2, 2.9, 1.85, 1.05, 17, INK, True)
tx(s, "timeout · auth · ratelimit은 이전 값 유지 / empty는 정상값 없이 unavailable", 1.1, 6.45, 11.1, 0.35, 14, MUTED, True)

# 09 Files
cards(
    "IMPLEMENTATION · 변경 파일과 역할",
    "계약·상태",
    ["widget.ts — DataStatus", "freshness.ts — 상태 투영", "clock.ts — 안정 1초 클록", "use-simulate.ts — URL 장애 입력"],
    "화면·검증",
    ["freshness-badge.tsx — 텍스트+아이콘", "dashboard.tsx — topbar·simulate", "t05-freshness.spec.ts — UI 검증", "vitest.config.ts — 실제 경로 별칭"],
    "docs/01_기획.md §6",
)

# 10 Protected
list_slide(
    "DO NOT TOUCH · 보호한 기존 경계",
    ["외부 API fetch·정규화: src/app/api/widgets/**", "기존 장애 재현 API: 재사용만", "기존 dashboard E2E 수정 없음", "고정검사·상한 문서 수정 없음", "기존 단위 테스트 assertion·기대값 변경 없음"],
    "docs/01_기획.md §7",
)

# 11 Checks overview
cards(
    "FIXED CHECKS · T05-C01~C10",
    "정상·timeout",
    ["C01 fresh", "C02 수신 시각·경과 초", "C03 이전 데이터 유지", "C04 stale", "C05 시각 동결·경과 증가"],
    "원인·접근성·자동화",
    ["C06 auth 401/403", "C07 ratelimit 대기", "C08 unavailable", "C09 텍스트 상태", "C10 E2E PASS"],
    "docs/00_과제_요구사항_매핑.md",
)

# 12-18 Evidence
evidence("C01·C02 · 정상 수신은 fresh", SHOT / "T05_01_fresh.png", "fresh · 마지막 수신 시각 · 경과 0s", "EVIDENCE · FRESH")
evidence("C03 · timeout에도 이전 데이터 유지", SHOT / "T05_02_stale_timeout.png", "마지막 정상값과 차트는 유지하고 배지만 stale로 전환", "EVIDENCE · TIMEOUT")
evidence("C04·C05 · stale 시각은 멈추고 경과는 증가", SHOT / "T05_02_stale_timeout.png", "마지막 정상 수신 시각 동결 · ageSeconds는 계속 증가", "EVIDENCE · STALE")
evidence("C06 · 권한 오류를 원인까지 표시", SHOT / "T05_03_stale_auth.png", "stale · 권한 오류(401/403)로 갱신 중지", "EVIDENCE · AUTH")
evidence("C07 · 호출 제한은 다음 갱신 대기", SHOT / "T05_04_stale_ratelimit.png", "stale · 호출 제한 — 다음 갱신 대기 중", "EVIDENCE · RATE LIMIT")
evidence("C08 · 정상값이 없으면 unavailable", SHOT / "T05_05_unavailable.png", "unavailable · 데이터를 가져오지 못했습니다", "EVIDENCE · EMPTY")
evidence("C09 · 색상 외 텍스트로 상태 구분", SHOT / "T05_01_fresh.png", "가시 텍스트와 aria-label로 fresh/stale/unavailable을 전달", "EVIDENCE · ACCESSIBILITY")

# 19 Test result
cards(
    "C10 · 전체 자동검사 결과",
    "품질 게이트",
    ["lint PASS", "typecheck PASS", "unit 45/45 PASS", "build PASS", "전체 E2E 14/14 PASS"],
    "최종 배포 재검증",
    ["공개 URL HTTP 200", "T05 원격 E2E 5/5", "5개 상태 육안 확인", "UTC 실행별 스크린샷 보존"],
    "docs/04_작업_기록.md",
)

# 20 Claude design
cards(
    "CLAUDE CODE · 설계 단계",
    "먼저 정한 것",
    ["관찰 가능한 고정검사 10개", "15분·사용자 프롬프트 10회 상한", "3값 dataStatus 계약", "수정 금지 경계"],
    "작업 분담",
    ["Codex는 정상 상태까지", "7항목 문서에서 멈춤", "Claude Code는 장애 분기·E2E", "대화 전문 없이 이어받기"],
    "docs/04_작업_기록.md · 설계",
)

# 21 Codex work
cards(
    "OPENAI CODEX · 1차 구현",
    "구현한 내용",
    ["DataStatus와 순수 투영 함수", "공용 1초 클록", "접근 가능한 FreshnessBadge", "dashboard topbar fresh 연동"],
    "종료 상태",
    ["정상 경로 화면 증거", "관련 테스트·lint·typecheck·build", "15분 · 사용자 프롬프트 1회", "장애 분기 전 인계"],
    "docs/02·04 · Primary Implementer",
)

# 22 Handoff
s = base("HANDOFF · 일곱 칸으로 멈추고 넘기기", "docs/02_인수인계_문서.md")
items = ["목표", "현재 상태", "실행 명령", "통과 검사", "남은 문제", "다음 행동", "건드리지 말 것"]
for i, item in enumerate(items):
    col = i % 4
    row = i // 4
    x = 0.65 + col * 3.05
    y = 1.65 + row * 2.35
    card(s, x, y, 2.75, 1.85)
    tx(s, f"{i + 1:02d}", x + 0.2, y + 0.25, 0.5, 0.3, 12, LIME, True)
    tx(s, item, x + 0.65, y + 0.62, 1.8, 0.45, 16, INK, True)

# 23 Claude follow-up
cards(
    "CLAUDE CODE · 인계 후 완성",
    "구현한 내용",
    ["simulate 4개 URL 분기", "이전 정상값 유지", "stale/unavailable 원인 매핑", "T05 E2E 5개"],
    "검증한 내용",
    ["고정검사 10/10", "전체 E2E 14/14", "14분 · 사용자 프롬프트 1회", "원격 상태 증거 5장"],
    "docs/04 · Secondary Implementer",
)

# 24-25 Corrections
cards(
    "CORRECTION 1 · 파일 누락이 아니었다",
    "처음 설명",
    ["fixture와 http 파일이 물리적으로 누락", "기존 테스트 import 실패"],
    "직접 확인한 실제 원인",
    ["파일은 저장소에 존재", "한글 경로가 URL 퍼센트 인코딩", "fileURLToPath로 실제 경로 사용", "단위 테스트 29/31 → 45/45"],
    "docs/02·04 · 정정 기록",
)
cards(
    "CORRECTION 2 · 화면 경계값과 누락 요구",
    "경과 초 표시",
    ["fresh 배지에 숫자 ageSeconds 누락", "경과 {n}s 문구 추가", "T05-C02 충족"],
    "음수 경계값",
    ["클라이언트 시계가 서버보다 뒤면 -1s", "ageSeconds 0 하한", "회귀 단위 테스트 추가", "전체 E2E 통과"],
    "docs/02·04 · 정정 기록",
)

# 26 Timeline content
s = base("WORK LOG · 실제 작업 내용으로 읽는 흐름", "docs/04_작업_기록.md")
stages = [
    ("Claude Code · 설계", "상태 규칙 3종\n· C01~C10 검사\n· 데이터 계약\n· 변경 금지 범위"),
    ("OpenAI Codex", "타입·상태 계산\n· 주입형 시계\n· 접근성 배지\n· fresh 화면"),
    ("인수인계 문서", "수정·생성 파일\n· 테스트 결과\n· 남은 작업\n· 금지 경계"),
    ("Claude Code · 완성", "simulate 장애 4종\n· stale/unavailable\n· 경계값 결함 수정\n· E2E 검증"),
    ("배포 · 최종 검증", "feature→main 통합\n· Vercel 공개\n· 5개 URL 육안 확인\n· 제출 문서 갱신"),
]
for i, (head, body) in enumerate(stages):
    x = 0.45 + i * 2.54
    card(s, x, 1.72, 2.25, 4.45)
    tx(s, head, x + 0.2, 2.08, 1.85, 0.3, 13, LIME, True)
    tx(s, body, x + 0.2, 2.82, 1.85, 2.25, 12, INK, True)
tx(s, "코드 번호가 아니라 각 단계의 결정·구현·검증 내용으로 재현", 2.05, 6.45, 9.3, 0.35, 14, MUTED, True)

# 27 Actual-name comparison
s = base("COMPARISON · 실제 역할명으로 비교", "docs/03·04 · 작업 측정")
heads = ["항목", "OpenAI Codex", "Claude Code"]
rows = [
    ("역할", "1차 구현", "설계 · 2차 구현"),
    ("구현 범위", "정상 계약·배지", "장애 분기·E2E"),
    ("작업 시간", "15분", "14분"),
    ("사용자 프롬프트", "1회", "1회"),
    ("최종 고정검사", "정상 범위 2/10", "전체 10/10"),
]
for col, head in enumerate(heads):
    x = [0.65, 4.48, 8.35][col]
    w = [3.58, 3.62, 4.33][col]
    card(s, x, 1.48, w, 0.68, CARD2)
    tx(s, head, x + 0.15, 1.72, w - 0.3, 0.25, 13, LIME, True)
for row_index, row in enumerate(rows):
    y = 2.32 + row_index * 0.75
    for col, value in enumerate(row):
        x = [0.65, 4.48, 8.35][col]
        w = [3.58, 3.62, 4.33][col]
        card(s, x, y, w, 0.62)
        tx(s, value, x + 0.15, y + 0.19, w - 0.3, 0.25, 12, INK, col == 0)

# 28 AI lines
list_slide(
    "AI 3줄 · 위임, 직접 판단, 불채택",
    ["AI에게 맡긴 일 — 신선도 배지, 이전 값 유지, simulate 분기, 자동 테스트와 인계 문서", "직접 판단한 일 — 에러 화면을 분리하지 않고 데이터 유지·배지만 3상태로 교체", "AI 말을 안 들은 일 — 파일 누락 진단을 직접 검증해 비ASCII 경로 인코딩을 실제 원인으로 확인"],
    "docs/AI_3줄.md",
    17,
)

# 29 Guide
s = base("VERIFICATION GUIDE · 30초, 3단계", "docs/검증안내서.md")
steps = [
    ("01 · /", "fresh · 수신 시각 · 경과 초"),
    ("02 · timeout", "이전 데이터 유지 · stale · 경과 증가"),
    ("03 · 나머지", "auth · ratelimit · empty"),
]
for i, (head, body) in enumerate(steps):
    x = 0.72 + i * 4.08
    card(s, x, 1.65, 3.78, 4.55)
    tx(s, head, x + 0.25, 2.0, 3.2, 0.35, 14, LIME, True)
    tx(s, body, x + 0.25, 2.8, 3.2, 1.3, 18, INK, True)
    tx(s, "텍스트 상태와\n안내 문구 확인", x + 0.25, 5.1, 3.1, 0.65, 12, MUTED)

# 30 Deploy
cards(
    "DEPLOYMENT · GitHub에서 Vercel 공개까지",
    "선택·설정",
    ["Vercel Hobby · Next.js", "Production Branch main", "Framework Preset Next.js", "Build·Output Auto-detect", "서버 키는 NEXT_PUBLIC_ 금지"],
    "최종 결과",
    ["GitHub PUBLIC", "무로그인 공개 URL HTTP 200", "원격 T05 E2E 5/5", "상태 5장 육안 확인", "UTC 폴더에 증거 누적"],
    "docs/05_배포.md",
)

# 31 Mapping
cards(
    "REQUIREMENTS · 완료와 남은 학생 확인",
    "완료",
    ["관찰형 고정검사 10/10", "공개 배포·GitHub PUBLIC", "시간·프롬프트·실패 회차 기록", "인계·검증안내·AI 3줄"],
    "추정하지 않고 남긴 항목",
    ["과제 원본 미제공 C32~C36", "과제 원본 미제공 C40~C49", "C37 개인정보 학생 점검", "C38 비밀값 학생 점검"],
    "docs/00_과제_요구사항_매핑.md",
)

# 32 Final
s = base("FINAL · 제출 주소와 결과", "T05 · FINAL LINKS")
card(s, 0.75, 1.55, 11.8, 4.85)
tx(s, "결과물", 1.05, 2.0, 1.4, 0.35, 13, LIME, True)
tx(s, "https://realtime-info-board-t05.vercel.app/", 2.55, 1.95, 9.0, 0.45, 18)
tx(s, "소스", 1.05, 3.05, 1.4, 0.35, 13, LIME, True)
tx(s, "https://github.com/stulss/realtime-info-board-t05", 2.55, 3.0, 9.0, 0.45, 18)
tx(s, "검증", 1.05, 4.1, 1.4, 0.35, 13, LIME, True)
tx(s, "고정검사 10/10 · unit 45/45 · E2E 14/14 · 원격 5/5", 2.55, 4.05, 9.0, 0.45, 17, INK, True)
tx(s, "원고", 1.05, 5.15, 1.4, 0.35, 13, LIME, True)
tx(s, "docs/T05_제출자료_원고_v3_2026-08-29.md", 2.55, 5.1, 9.0, 0.45, 16)

assert len(prs.slides) == 32
prs.save(OUT)
print(OUT)
