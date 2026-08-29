from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SHOT_DIR = ROOT / "docs" / "검증스크린샷" / "2026-08-29T04-21-36Z"
OUT = ROOT / "docs" / "T05_과제_제출_보고서_통합본_2026-08-29.pptx"
if OUT.exists():
    raise SystemExit(f"refusing to overwrite {OUT}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(248, 247, 241)
CARD = RGBColor(255, 255, 255)
INK = RGBColor(22, 25, 21)
MUTED = RGBColor(93, 98, 88)
LIME = RGBColor(198, 255, 34)
GREEN = RGBColor(40, 128, 77)
AMBER = RGBColor(184, 113, 20)
RED = RGBColor(185, 47, 37)


def textbox(slide, text, x, y, w, h, size=17, color=INK, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return box


def card(slide, x, y, w, h):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(220, 220, 210)
    return shape


def base(title, kicker="T05 · PULSEBOARD"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.58), Inches(0.22), Inches(2.55), Inches(0.34)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = LIME
    tag.line.fill.background()
    textbox(slide, kicker, 0.7, 0.27, 2.25, 0.2, 8, INK, True)
    textbox(slide, title, 0.58, 0.72, 11.7, 0.65, 25, INK, True)
    textbox(slide, f"{len(prs.slides):02d}", 12.2, 7.05, 0.5, 0.2, 8, MUTED)
    return slide


def bullets(slide, items, x, y, w, h, size=17, color=INK):
    textbox(slide, "\n".join(f"• {item}" for item in items), x, y, w, h, size, color)


def place_image(slide, path, x, y, w, h):
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    slide.shapes.add_picture(
        str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph)
    )


def evidence(slide, path, x, y, w, h, label, color):
    card(slide, x, y, w, h)
    textbox(slide, label, x + 0.2, y + 0.18, w - 0.4, 0.3, 12, color, True)
    place_image(slide, path, x + 0.16, y + 0.58, w - 0.32, h - 0.75)


# 1. Cover
s = base("데이터가 늦어도 화면은 사라지지 않는다", "T05 · 과제 제출 보고서")
textbox(s, "실시간 정보판 데이터 신선도 배지", 0.72, 1.55, 7.3, 0.55, 26, INK, True)
textbox(s, "fresh · stale · unavailable\n마지막 정상값 유지 + 경과 시간 표시", 0.72, 2.35, 7.4, 1.15, 20, MUTED)
card(s, 0.72, 4.05, 11.8, 1.5)
textbox(s, "결과물  https://realtime-info-board-t05.vercel.app/\n소스     https://github.com/stulss/realtime-info-board-t05", 1.02, 4.42, 11.1, 0.8, 16, INK)

# 2. Goal and solution
s = base("01 기획 · 한 화면에서 데이터 현재성 판단", "docs/01_기획.md")
card(s, 0.65, 1.55, 5.9, 4.95)
textbox(s, "문제", 0.95, 1.88, 5.2, 0.35, 14, RED, True)
bullets(s, ["장애 시 기존 데이터까지 사라짐", "마지막 정상 수신 시각을 알기 어려움", "색상만으로 상태를 구분하면 접근성 저하"], 0.95, 2.45, 5.1, 2.5)
card(s, 6.78, 1.55, 5.9, 4.95)
textbox(s, "해결", 7.08, 1.88, 5.2, 0.35, 14, GREEN, True)
bullets(s, ["fresh / stale / unavailable 3상태", "장애 시 마지막 정상값·차트 유지", "시각·경과 초·텍스트·aria-label 제공"], 7.08, 2.45, 5.1, 2.5)

# 3. AI 3 lines
s = base("AI 3줄 · 위임, 직접 판단, 불채택", "docs/AI_3줄.md")
card(s, 0.65, 1.5, 12.03, 1.35)
textbox(s, "AI에게 맡긴 일", 0.95, 1.82, 2.0, 0.3, 13, GREEN, True)
textbox(s, "신선도 배지 · 이전 값 유지 · simulate 분기 · 자동 테스트 · 인계 문서", 3.0, 1.78, 9.1, 0.5, 16, INK)
card(s, 0.65, 3.0, 12.03, 1.35)
textbox(s, "내가 판단한 일", 0.95, 3.32, 2.0, 0.3, 13, AMBER, True)
textbox(s, "에러 화면을 분리하지 않고 데이터는 유지한 채 배지만 3상태로 교체", 3.0, 3.28, 9.1, 0.5, 16, INK)
card(s, 0.65, 4.5, 12.03, 1.35)
textbox(s, "AI 말을 안 들은 일", 0.95, 4.82, 2.0, 0.3, 13, RED, True)
textbox(s, "파일 누락 진단을 직접 검증해 실제 원인인 비ASCII Vitest 경로 인코딩을 수정", 3.0, 4.78, 9.1, 0.55, 16, INK)

# 4. Requirements mapping
s = base("00 요구사항 매핑 · 고정검사와 루브릭", "docs/00_과제_요구사항_매핑.md")
card(s, 0.65, 1.5, 5.75, 5.05)
textbox(s, "고정검사 T05-C01~C10", 0.95, 1.85, 5.1, 0.35, 14, GREEN, True)
bullets(s, ["fresh·시각·경과 초", "timeout 이전 데이터 유지", "auth·ratelimit 원인 문구", "empty unavailable", "텍스트 접근성·E2E"], 0.95, 2.35, 5.0, 3.25, 16)
textbox(s, "10 / 10 PASS", 1.0, 5.85, 4.7, 0.35, 19, GREEN, True)
card(s, 6.65, 1.5, 6.03, 5.05)
textbox(s, "과제 5 루브릭", 6.95, 1.85, 5.4, 0.35, 14, GREEN, True)
bullets(s, ["사전 고정·공통 상한 완료", "AI A 중단·7항목 인계 완료", "AI B 새 대화 이어받기 완료", "익명 비교·공개 배포 완료", "원본 미제공 C32~36·C40~49만 대조 대기"], 6.95, 2.35, 5.25, 3.25, 16)
textbox(s, "C21 공개·GitHub PUBLIC ✅", 7.0, 5.85, 5.0, 0.35, 17, GREEN, True)

# 5. Handoff
s = base("02 인수인계 · AI A에서 AI B로", "docs/02_인수인계_문서.md")
card(s, 0.65, 1.5, 3.65, 5.1)
textbox(s, "AI A", 0.95, 1.85, 3.0, 0.35, 14, GREEN, True)
bullets(s, ["타입·순수 함수", "공용 1초 클록", "접근 가능한 배지", "fresh 렌더", "15분 · 1회 · 2/10"], 0.95, 2.35, 3.0, 3.2, 16)
card(s, 4.58, 1.5, 4.15, 5.1)
textbox(s, "7항목 인계", 4.88, 1.85, 3.5, 0.35, 14, AMBER, True)
bullets(s, ["목표·현재 상태", "실행 명령·통과 검사", "남은 문제·다음 행동", "건드리지 말 것", "인계 커밋 90be0ac"], 4.88, 2.35, 3.45, 3.2, 16)
card(s, 9.0, 1.5, 3.68, 5.1)
textbox(s, "AI B", 9.3, 1.85, 3.0, 0.35, 14, GREEN, True)
bullets(s, ["simulate 4분기", "stale/unavailable", "경계 버그 수정", "E2E 5건", "14분 · 1회 · 10/10"], 9.3, 2.35, 3.0, 3.2, 16)

# 6. Work timeline
s = base("04 작업 기록 · 설계에서 공개까지", "docs/04_작업_기록.md")
stages = [
    ("설계", "fa41377", "검사·상한·문서 동결"),
    ("AI A", "90be0ac", "15분 · 1회 · 2/10"),
    ("AI B", "dd87a46", "14분 · 1회 · 10/10"),
    ("통합", "326c7f5", "main 머지 · PUBLIC"),
    ("마감", "33c76c8", "배포 · 증거 · PPT/PDF"),
]
for i, (stage, commit, detail) in enumerate(stages):
    x = 0.58 + i * 2.52
    card(s, x, 1.75, 2.28, 4.45)
    textbox(s, stage, x + 0.2, 2.08, 1.85, 0.35, 14, GREEN, True)
    textbox(s, commit, x + 0.2, 2.75, 1.85, 0.35, 13, INK, True)
    textbox(s, detail, x + 0.2, 3.55, 1.85, 1.15, 15, INK)
textbox(s, "AI B 실패 회차 2회 → 경계값 수정 후 전체 E2E 14/14 PASS", 1.25, 6.48, 10.8, 0.35, 15, MUTED, True)

# 7. Architecture
s = base("01 기획 · 기존 경계를 보존한 최소 변경", "docs/01_기획.md")
for i, (head, body) in enumerate([
    ("URL", "?simulate=\n4개 장애 입력"),
    ("QUERY", "기존 fetch와\n실패 처리 재사용"),
    ("STATE", "6값 status →\n3값 dataStatus"),
    ("VIEW", "한 화면 유지\n배지만 교체"),
]):
    x = 0.72 + i * 3.08
    card(s, x, 1.75, 2.75, 3.7)
    textbox(s, head, x + 0.25, 2.1, 2.25, 0.35, 14, GREEN, True)
    textbox(s, body, x + 0.25, 2.85, 2.25, 1.25, 18, INK, True)
textbox(s, "외부 API·DB·기존 E2E 기대값 무수정", 2.7, 6.05, 8.0, 0.4, 16, MUTED, True)

# 5. Fresh / timeout
s = base("EVIDENCE · fresh에서 stale로")
evidence(s, SHOT_DIR / "T05_01_fresh.png", 0.58, 1.5, 6.05, 4.95, "fresh · 정상 수신 · 경과 0s", GREEN)
evidence(s, SHOT_DIR / "T05_02_stale_timeout.png", 6.72, 1.5, 6.05, 4.95, "stale · timeout · 이전 데이터 유지", AMBER)

# 6. Other failures
s = base("EVIDENCE · 원인별 stale와 unavailable")
evidence(s, SHOT_DIR / "T05_03_stale_auth.png", 0.45, 1.55, 4.05, 4.8, "stale · 권한 오류 401/403", AMBER)
evidence(s, SHOT_DIR / "T05_04_stale_ratelimit.png", 4.64, 1.55, 4.05, 4.8, "stale · 호출 제한", AMBER)
evidence(s, SHOT_DIR / "T05_05_unavailable.png", 8.83, 1.55, 4.05, 4.8, "unavailable · 정상값 없음", RED)

# 10. Verification guide
s = base("검증안내서 · 30초, 3단계", "docs/검증안내서.md")
steps = [
    ("01 · /", "fresh · 수신 시각 · 경과 초", "C01·C02"),
    ("02 · timeout", "이전 데이터 유지 · stale · 경과 증가", "C03~C05"),
    ("03 · 나머지", "auth · ratelimit · empty", "C06~C08"),
]
for i, (head, body, ids) in enumerate(steps):
    x = 0.7 + i * 4.15
    card(s, x, 1.65, 3.8, 4.55)
    textbox(s, head, x + 0.25, 2.0, 3.25, 0.35, 14, GREEN, True)
    textbox(s, body, x + 0.25, 2.75, 3.25, 1.4, 18, INK, True)
    textbox(s, ids, x + 0.25, 5.25, 3.25, 0.35, 14, MUTED, True)
textbox(s, "색상뿐 아니라 텍스트·aria-label로 구분(C09) · E2E 전체 통과(C10)", 1.45, 6.5, 10.4, 0.35, 15, MUTED, True)

# 11. Verification
s = base("VERIFICATION · 고정검사 10/10")
card(s, 0.65, 1.55, 5.5, 4.95)
textbox(s, "품질 게이트", 0.95, 1.9, 4.8, 0.35, 14, GREEN, True)
bullets(s, ["lint PASS", "typecheck PASS", "unit 45/45 PASS", "build PASS", "E2E 14/14 PASS"], 0.95, 2.4, 4.7, 3.2, 17)
card(s, 6.45, 1.55, 6.23, 4.95)
textbox(s, "최종 배포 확인", 6.75, 1.9, 5.5, 0.35, 14, GREEN, True)
bullets(s, ["공개 URL 로그인 없음", "원격 T05 E2E 5/5 PASS", "fresh → stale → unavailable 육안 확인", "증거: 2026-08-29T04-21-36Z"], 6.75, 2.4, 5.35, 3.2, 17)

# 13. Deployment
s = base("05 배포 · GitHub에서 Vercel 공개까지", "docs/05_배포.md")
card(s, 0.65, 1.5, 5.85, 5.05)
textbox(s, "선택과 설정", 0.95, 1.85, 5.2, 0.35, 14, GREEN, True)
bullets(s, ["Vercel Hobby · Next.js 단일 스택", "Production Branch = main", "Framework Preset = Next.js", "Build/Output = Auto-detect", "서버 키는 NEXT_PUBLIC_ 금지"], 0.95, 2.35, 5.05, 3.3, 16)
card(s, 6.78, 1.5, 5.9, 5.05)
textbox(s, "최종 결과", 7.08, 1.85, 5.2, 0.35, 14, GREEN, True)
bullets(s, ["GitHub 저장소 PUBLIC", "무로그인 공개 URL HTTP 200", "원격 T05 E2E 5/5", "fresh→stale→unavailable 육안 확인", "증거 2026-08-29T04-21-36Z"], 7.08, 2.35, 5.05, 3.3, 16)

# 14. Comparison
s = base("COMPARISON · 이름을 가린 AI A/B")
headers = ["항목", "AI A", "AI B"]
rows = [
    ("범위", "fresh 계약·UI", "장애 분기·E2E"),
    ("시간", "15분", "14분"),
    ("사용자 프롬프트", "1회", "1회"),
    ("실패 회차", "0회", "2회"),
    ("최종 고정검사", "2/10", "10/10"),
]
for col, head in enumerate(headers):
    x = [0.7, 4.55, 8.5][col]
    w = [3.6, 3.65, 4.1][col]
    card(s, x, 1.55, w, 0.65)
    textbox(s, head, x + 0.15, 1.75, w - 0.3, 0.25, 13, GREEN, True)
for row_index, row in enumerate(rows):
    y = 2.35 + row_index * 0.72
    for col, value in enumerate(row):
        x = [0.7, 4.55, 8.5][col]
        w = [3.6, 3.65, 4.1][col]
        card(s, x, y, w, 0.58)
        textbox(s, value, x + 0.15, y + 0.17, w - 0.3, 0.22, 12, INK, col == 0)

# 15. Problems and judgment
s = base("LESSONS · 기록을 믿되 반드시 검증한다")
card(s, 0.65, 1.55, 5.9, 4.95)
textbox(s, "발견한 문제", 0.95, 1.9, 5.2, 0.35, 14, RED, True)
bullets(s, ["인계 문서의 파일 누락 오진", "한글 경로의 Vitest 별칭 인코딩", "경과 -1s 경계값", "Vercel Other 프리셋의 dist 오류"], 0.95, 2.4, 5.1, 3.2, 16)
card(s, 6.78, 1.55, 5.9, 4.95)
textbox(s, "해결 원칙", 7.08, 1.9, 5.2, 0.35, 14, GREEN, True)
bullets(s, ["원문보다 저장소 사실을 우선", "fileURLToPath로 실제 경로 사용", "ageSeconds 0 하한 + 회귀 테스트", "Next.js 프리셋·자동 감지"], 7.08, 2.4, 5.1, 3.2, 16)

# 16. Final links
s = base("FINAL · 제출 주소와 결과")
card(s, 0.72, 1.55, 11.85, 4.95)
textbox(s, "결과물", 1.05, 2.0, 1.4, 0.3, 13, GREEN, True)
textbox(s, "https://realtime-info-board-t05.vercel.app/", 2.55, 1.95, 9.2, 0.42, 18, INK)
textbox(s, "소스", 1.05, 3.0, 1.4, 0.3, 13, GREEN, True)
textbox(s, "https://github.com/stulss/realtime-info-board-t05", 2.55, 2.95, 9.2, 0.42, 18, INK)
textbox(s, "결과", 1.05, 4.0, 1.4, 0.3, 13, GREEN, True)
textbox(s, "고정검사 10/10 · unit 45/45 · E2E 14/14 · 원격 5/5", 2.55, 3.95, 9.2, 0.42, 18, INK, True)
textbox(s, "다음 기준", 1.05, 5.0, 1.4, 0.3, 13, GREEN, True)
textbox(s, "초기 계약은 기록 중심, 장애·경계는 실행·검증 루프 중심", 2.55, 4.95, 9.2, 0.42, 16, MUTED)

assert len(prs.slides) == 15
prs.save(OUT)
print(OUT)
